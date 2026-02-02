#!/usr/bin/env python3
"""Создание PPTX презентации для WMS Buffer Optimization"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title box
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, note=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)

    # Content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = bullet
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.level = 0

    if note:
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(128, 128, 128)

    return slide

def add_diagram_slide(prs, title, diagram_text, bullets=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)

    # Diagram box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    if bullets:
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(2.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            p.space_before = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)

    # Table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5)

    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Оптимизация буферной зоны склада",
        "WMS Buffer Management System")

    # Slide 2: Current situation
    add_diagram_slide(prs,
        "1. Текущая ситуация",
        "ХРАНЕНИЕ  ──[3 карщика]──►  БУФЕР (64 ячейки)  ──[20 сборщиков]──►  СБОРКА",
        [
            "• 3 карщика доставляют моно-палеты из хранения в буфер",
            "• 64 ячейки буфера — временное хранение для сборки",
            "• 20 сборщиков забирают товары и распределяют по заказам",
            "• Проблема: синхронизация подачи палет со скоростью сборки"
        ])

    # Slide 3: Why it matters
    add_table_slide(prs,
        "2. Почему это важно — потери",
        ["Ситуация", "Последствия"],
        [
            ["Буфер пустой (<15%)", "Сборщики простаивают, заказы задерживаются"],
            ["Буфер переполнен (>70%)", "Карщики ждут освобождения ячеек"],
            ["Неравномерная нагрузка", "Одни перегружены, другие простаивают"],
            ["Каждый сбой потока", "Задержка целой волны заказов"]
        ])

    # Slide 4: Solution overview
    add_content_slide(prs,
        "3. Решение — 3 уровня оптимизации",
        [
            "УРОВЕНЬ 3: ПРОГНОЗИРОВАНИЕ (Historical Layer)",
            "   → ML-модели предсказывают время выполнения задач",
            "",
            "УРОВЕНЬ 2: ПЛАНИРОВАНИЕ (Tactical Layer)",
            "   → OR-Tools оптимизирует назначения и расписание",
            "",
            "УРОВЕНЬ 1: РЕАГИРОВАНИЕ (Realtime Layer)",
            "   → Гистерезис-контроллер управляет уровнем буфера",
            "",
            "Принцип: Прогноз → План → Исполнение → Корректировка"
        ])

    # Slide 5: ML Models
    add_content_slide(prs,
        "4. ML-модели для прогнозирования",
        [
            "Модель 1: Время сборки (Picker Duration)",
            "   Входы: ID сборщика, кол-во строк, кол-во товаров, час, день недели",
            "   Выход: прогноз времени выполнения задания (секунды)",
            "",
            "Модель 2: Время маршрута карщика (Forklift Duration)",
            "   Входы: ID карщика, зона источника → зона назначения, вес палеты",
            "   Выход: прогноз времени доставки (секунды)",
            "",
            "База для обучения: 1.5 млн исторических записей"
        ])

    # Slide 6: Smart Picker Assignment
    add_content_slide(prs,
        "5. Умное назначение задач пикерам",
        [
            "Проблема: не все пикеры одинаково эффективны для всех товаров",
            "",
            "Решение — назначаем задачу тому, кто:",
            "",
            "   • Лучше справляется с данным типом товара (из истории)",
            "   • Сейчас свободен или скоро освободится",
            "   • Находится ближе к нужной зоне буфера",
            "   • Имеет меньшую загрузку в текущей волне",
            "",
            "ML модель учитывает связку: пикер + товар + объём → время"
        ])

    # Slide 7: Product Statistics
    add_content_slide(prs,
        "6. Аналитика по товарам",
        [
            "Для каждого товара рассчитываем:",
            "",
            "   • Среднее время распределения",
            "   • Вариативность времени (стабильность)",
            "   • Типичное количество в задании",
            "   • Частота появления",
            "",
            "Классификация сложности товара (1-10):",
            "   complexity = 0.4×время + 0.3×вариативность + 0.2×qty + 0.1×редкость",
            "",
            "   1-3: лёгкие товары → любой пикер",
            "   4-6: средние → стандартное назначение",
            "   7-10: сложные → опытным пикерам"
        ])

    # Slide 8: OR-Tools
    add_content_slide(prs,
        "7. Оптимизация расписания (OR-Tools)",
        [
            "Задача: минимизация общего времени выполнения волны",
            "",
            "Оптимизируем ДВА назначения одновременно:",
            "   • Карщик → Палета → Ячейка буфера",
            "   • Пикер → Задание на сборку",
            "",
            "Учитываем:",
            "   • Прогноз времени для каждой пары (пикер, задание)",
            "   • Текущую занятость и время до освобождения",
            "   • Балансировку нагрузки между всеми работниками",
            "",
            "Ограничения: приоритеты заказов, тяжёлые палеты вниз"
        ])

    # Slide 9: How it works
    add_diagram_slide(prs,
        "8. Цикл оптимизации (каждые 15 минут)",
        "Заказы (волна) → Прогноз ML.NET → OR-Tools CP-SAT → Исполнение WMS → Обратная связь",
        [
            "1. Получаем список заказов для сборки",
            "2. ML предсказывает время каждой операции",
            "3. OR-Tools строит оптимальный план",
            "4. План отправляется в WMS для исполнения",
            "5. Результаты используются для улучшения моделей"
        ])

    # Slide 10: Example Forklifts
    add_content_slide(prs,
        "9. Пример: карщики (волна 50 заданий)",
        [
            "БЕЗ оптимизации:",
            "   Карщик 1: ████████████████████████ (перегружен)",
            "   Карщик 2: ████████░░░░░░░░░░░░░░░░ (недогружен)",
            "   Карщик 3: ██████████████░░░░░░░░░░ (средне)",
            "",
            "С оптимизацией OR-Tools:",
            "   Карщик 1: ████████████████░░░░░░░░ (сбалансировано)",
            "   Карщик 2: ███████████████░░░░░░░░░ (сбалансировано)",
            "   Карщик 3: ████████████████░░░░░░░░ (сбалансировано)",
            "",
            "Результат: время волны 45→32 мин (-29%)"
        ])

    # Slide 11: Example Pickers
    add_content_slide(prs,
        "10. Пример: пикеры (назначение по эффективности)",
        [
            "БЕЗ оптимизации (случайное назначение):",
            "   Пикер А берёт товар X → 8 мин (не его специализация)",
            "   Пикер Б берёт товар Y → 10 мин (не его специализация)",
            "",
            "С оптимизацией (ML + OR-Tools):",
            "   Пикер А берёт товар Y → 5 мин (его лучший товар)",
            "   Пикер Б берёт товар X → 4 мин (его лучший товар)",
            "",
            "Учитываем: кто скоро освободится, кто ближе к ячейке",
            "",
            "Результат: простои пикеров 12%→3% (-75%)"
        ])

    # Slide 12: Expected results
    add_table_slide(prs,
        "11. Ожидаемые результаты",
        ["Метрика", "До", "После", "Улучшение"],
        [
            ["Время волны", "45 мин", "32 мин", "-29%"],
            ["Простои сборщиков", "12%", "3%", "-75%"],
            ["Критические ситуации", "5/день", "<1/день", "-80%"],
            ["Пробег карщиков", "100%", "85%", "-15%"]
        ])

    # Slide 13: Implementation plan
    add_table_slide(prs,
        "12. План внедрения",
        ["Этап", "Задачи", "Результат"],
        [
            ["1. ML-модели", "Обучение на истории", "Прогноз времени"],
            ["2. Тестирование", "Сравнение с реальностью", "Оценка точности"],
            ["3. OR-Tools", "Интеграция оптимизатора", "Авто-планирование"],
            ["4. Пилот", "Работа на реальных данных", "Валидация"],
            ["5. Продакшн", "Полная интеграция с WMS", "Автономная работа"]
        ])

    # Slide 14: Summary
    add_content_slide(prs,
        "13. Резюме",
        [
            "Система WMS Buffer Management позволит:",
            "",
            "• Назначать задачи лучшим исполнителям (пикер↔товар)",
            "• Учитывать занятость и время освобождения",
            "• Оптимизировать маршруты карщиков",
            "• Сократить время волны на ~30%",
            "• Минимизировать простои персонала",
            "",
            "Основа: 1.5 млн записей + ML.NET + OR-Tools"
        ])

    # Save
    output_path = "/home/onica_on/WMS.BufferManagement/docs/WMS_Buffer_Optimization.pptx"
    prs.save(output_path)
    print(f"Презентация сохранена: {output_path}")

if __name__ == "__main__":
    main()
